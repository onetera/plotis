# :coding: utf-8

import re
from core import Core
from db_conn import DBconn
from synop import Synop

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class PPT( Core ):
    def write_ppt( self, scenario, scenario_idx ):
        search_msg  = '이 시스템은 영화 마케팅 전문가 입니다.'
        search_msg += '\n'
        search_msg += '아래는 제작을 하려고 하는 시나리오 입니다'
        search_msg += '{body}'
        search_msg += '\n'
        search_msg += '위 시나리오를 바탕으로 제작한 영화에 투자자들의 투자를 결정할 수 있는'
        search_msg += '프레젠테이션용 슬라이드 초안을 작성해주세요.'
        search_msg += '슬라이드 주제와 포함 정보는 아래와 같다.'
        search_msg += '\n'
        search_msg += '슬라이드 구성안'
        search_msg += '슬라이드 1: 영화 소개 - 제목, 캐치프레이즈, 장르, 예상 러닝타임, 타겟목표 관객층, 예상 관람등급 등'
        search_msg += '슬라이드 2 : 기획의도 (서정적이고 감성적인 문체를 사용해 900 ~ 1000자 내외 작성),'
        search_msg += '슬라이드 3: 시놉시스,'
        search_msg += '슬라이드 4: 관전포인트 (4가지 이상, 각 포인트는 300~400자로 서술해 감정을 자극할 것),'
        search_msg += '슬라이드 5: 캐릭터 소개 (주요 캐릭터 4명 이상 소개, 각 캐릭터의 특징이 드러나는 대사를 활용해 400~500자 서술)'
        search_msg += '\n'
        search_msg += '요구된 항목 외에 다른 키워드는 사용하지 말고 위의 가이드를 따라 작성'
        search_msg += '\n'
        search_msg += '작성형식)### [슬라이드 1: 영화 소개]'
        search_msg += '- 제목: 영화 제목'
        search_msg += '### [슬라이드 2: 기획의도]'
        search_msg += '기획의도의 내용'
        search_msg += '### [슬라이드 3: 시놉시스]'
        search_msg += '시놉시스의 내용'
        search_msg += '### [슬라이드 4: 관전포인트]'
        search_msg += '1. 첫번째 관전포인트 내용의 명사형'
        
        chain = self.chain( search_msg )
        
        texts = chain.invoke( {'body':scenario} )

        texts_dic = self.parse_ppt(texts)

        ppt_path = self.make_ppt_file(texts_dic, scenario, scenario_idx)

        return ppt_path
    
    def parse_ppt(self, texts):
        title_pattern = re.compile(r'\#\#\# \[슬라이드 \d+: (.*?)\]')
        content_pattern = re.compile(r'\]\n*(.*?)(?:\#\#\#|\Z)', re.DOTALL)

        title_lists = title_pattern.findall(texts)

        texts_dic = {}

        for i in range(len(title_lists)):
            if i != len(title_lists)-1:
                text_tmp = re.search(title_lists[i+1], texts)
                text = texts[:(text_tmp.end()+1)]

                content_lists = content_pattern.findall(text)

                texts = texts[text_tmp.end():]
            else:
                content_lists = content_pattern.findall(texts)

            texts_dic[title_lists[i]] = content_lists[0].strip()
        
        return texts_dic
    
    def make_ppt_file(self, texts_dic, scenario, scenario_idx):
        prs = Presentation( './tmp/template.pptx' )

        template = prs.slides[0]

        for i, (slide_title, slide_content) in enumerate(texts_dic.items()):
            if i < len(prs.slides):
                slide = prs.slides[i]
            else:
                slide_layout = template.slide_layout
                slide = prs.slides.add_slide(slide_layout)

                template_title_shape = template.shapes.title
                title_shape = slide.shapes.title
                self.title_format(template_title_shape, title_shape)

            title = slide.shapes.title
            title.text = slide_title

            left = title.left
            top = title.top + title.height + Inches(0.2)
            width = Inches(9) 
            height = Inches(4) 

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            if slide_title == '시놉시스':
                created = self.db.search_created(scenario_idx)
                if created:
                    p = text_frame.add_paragraph()
                    p.text = self.db.last_synop()[0][1] 
                    p.font.size = Pt(14) 
                    p.alignment = PP_ALIGN.LEFT
                else:
                    p = text_frame.add_paragraph()
                    synop = Synop()
                    p.text = synop.analyze_synop(scenario)
                    p.font.size = Pt(14) 
                    p.alignment = PP_ALIGN.LEFT
            
            else:
                sub_content = slide_content.replace('**', '').split('\n')

                for k in sub_content:
                    p = text_frame.add_paragraph()
                    p.text = k.strip()
                    p.font.size = Pt(14) 
                    p.alignment = PP_ALIGN.LEFT

        ppt_path = './tmp/proposal.pptx'
        prs.save(ppt_path)
        search_ppt = self.db.load_ppt_path(scenario_idx)
        if search_ppt:
            self.db.update_ppt(ppt_path, scenario_idx)
        else:
            self.db.insert_ppt( ppt_path, scenario_idx )
        
        return ppt_path

    def title_format(self, template_shape, new_shape):
        new_shape.left = template_shape.left
        new_shape.top = template_shape.top
        new_shape.width = template_shape.width
        new_shape.height = template_shape.height
